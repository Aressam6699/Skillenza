from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import OrderedDict
import re
import os
n=0

def remove_clean(f):
    escapes = ''.join([chr(char) for char in range(1, 32)])
    translator = str.maketrans('', '', escapes)
    f = f.translate(translator)
    f=f.strip()
    if len(f)==0:
        f=''
    return f
    
    

def write_image(shape):
    global n
    image = shape.image
    # ---get image "file" contents---
    image_bytes = image.blob
    # ---make up a name for the file, e.g. 'image.jpg'---
    image_filename = 'image{:03d}.{}'.format(n, image.ext)
    n += 1
    #print(image_filename)
    imname=image_filename
    with open('templates/'+image_filename, 'wb') as f:
        f.write(image_bytes)
    return imname

def visitor(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            visitor(s)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        imname=write_image(shape)
        return imname




def parseppts(filename):
    urls=[]
    
    prs = Presentation("upload/"+filename)

    pptcontent=OrderedDict() 
    #{slideid:[titles[],subtitles[],text[],pictures[],tables[]]}
    
    n=0

    idd=0
    for slide in prs.slides:
        #print(slide.slide_id)
        titles=[]
        subtitles=[]
        text=[]
        tables=[]
        images=[]
        for shape in slide.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if 'title' in str(phf.type).lower():
                    if 'subtitle' in str(phf.type).lower():
                        subtitles.append(remove_clean(shape.text))
                    else:
                        y=remove_clean(shape.text)
                        if len(y)>1:
                            titles=[y]
                        
            
            
            if len(titles)==0:
                #print(pptcontent[idd-1][0][0])
                try:
                    titles.append(pptcontent[idd-1][0][0])
                except:

                    titles=['NEW PPT']
                #print('+++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++')
                #break
            else:
                aaa=titles[0]
            
                if 'contd' in aaa.lower():
                    try:
                        titles=[pptcontent[idd-1][0][0]]
                    except:
                        titles=['NEW PPT']

                    #print('yeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeeet')
                #titles.append(pptcontent[idd][0])
                
            for shape in slide.shapes:
                imname=visitor(shape)
                cur_path=os.getcwd()
                os.chdir('templates')
                if imname:
                    if len(images)==0:
                        images.append(imname)
                    else:
                        exist=False
                        for j in images:
                            if open(j,"rb").read() == open(imname,"rb").read():
                                exist=True
                                break
                        if not exist:
                            images.append(imname)
                os.chdir(cur_path)

            



            for shape in slide.shapes:

                if shape.has_table:
                    text_runs=[]

                    tbl = shape.table
                    row_count = len(tbl.rows)
                    col_count = len(tbl.columns)
                    for r in range(0, row_count):
                        for c in range(0, col_count):
                            cell = tbl.cell(r,c)
                            paragraphs = cell.text_frame.paragraphs
                            cell_t=''
                            for paragraph in paragraphs:
                                for run in paragraph.runs:
                                    cell_t+=run.text
                            text_runs.append(run.text)
                    tab=[]
                    #print(len(text_runs),row_count*col_count)
                    try:
                        for i in range(row_count):
                            rr=[]
                            for j in range(col_count):
                                rr.append(text_runs[0])
                                text_runs.pop(0)
                            tab.append(rr)
                        if tab not in tables:
                            tables.append(tab)
                    except Exception as e:
                        
                        print(tab)
                        print(e)
                    
            for shape in slide.shapes:
                
                if hasattr(shape,"text"):
                    if not shape.has_table:
                        if shape.has_text_frame:
                                #print(shape.shape_type)

                                for paragraph in shape.text_frame.paragraphs:
                                    cleansed=remove_clean(paragraph.text)
                                    u = re.findall(r'^https?:\/\/.*[\r\n]*', cleansed)
                                    urls+=u
                                    text.append(cleansed)

                        else:

                            cleansed=remove_clean(shape.text)
                            u = re.findall(r'^https?:\/\/.*[\r\n]*', cleansed)
                            urls+=u
                            if cleansed not in text:
                                text.append(cleansed)
                        '''cleansed=remove_clean(shape.text)
                        u = re.findall(r'^https?:\/\/.*[\r\n]*', cleansed)
                        urls+=u
                        text.append(cleansed)
                            
                       if shape.is_placeholder:
                            
                            phf = shape.placeholder_format
                            if 'title' not in str(phf.type).lower():
                                if not shape.has_text_frame:
                                    cleansed=remove_clean(shape.text)
                                    u = re.findall(r'^https?:\/\/.*[\r\n]*', cleansed)
                                    urls+=u
                                    if cleansed not in text:
                                        text.append(cleansed)
                            else:'''

                            

                                

        #print(len(text))
        if len(text)!=0:
            copy=text.pop(0)
            temptext=[]
            for i in text:
                if i !=copy:
                    if i not in temptext:
                        temptext.append(i)
            text=temptext
        
        pptcontent.update({idd:[titles,subtitles,text,images,tables]})
        idd+=1
    return pptcontent,urls

def title(pptcontent):
    ppttitleid=min(pptcontent.keys())
    #print(ppttitleid)
    ppttitle=pptcontent[ppttitleid][0][0]
    return ppttitle
    #print(ppttitle)
        
                            

def readppts(directory):
    
    contents=OrderedDict()
    links={}

    
    for i in directory:
        x,url=parseppts(i)
        print(i)
        name=title(x)
        links.update({name:list(set(url))})
        contents.update({name:x})
    return contents,links


'''directory=['1.pptx','15.pptx','16.pptx']
readppts(directory)'''
                
            
      